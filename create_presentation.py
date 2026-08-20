from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

green = RGBColor(18, 127, 77)
dark = RGBColor(22, 31, 28)
light = RGBColor(243, 247, 238)
white = RGBColor(255, 255, 255)


def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = light

    banner = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.35))
    banner.fill.solid()
    banner.fill.fore_color.rgb = green
    banner.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(8.5), Inches(0.7))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = white
    p.alignment = PP_ALIGN.LEFT

    subtitle_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(9.5), Inches(0.5))
    p = subtitle_box.text_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(230, 245, 236)

    cards = [
        (Inches(0.7), Inches(1.9), Inches(3.75), Inches(2.45), 'Problem Statement\n\nAgriculture in Kenya faces uncertainty in estimating harvest output. Traditional methods are slow and limited.'),
        (Inches(4.8), Inches(1.9), Inches(3.75), Inches(2.45), 'Objective\n\nBuild a machine learning model to predict crop yield using remote sensing, soil, and weather data.'),
        (Inches(8.9), Inches(1.9), Inches(3.75), Inches(2.45), 'Methodology\n\nFeature engineering, encoding, scaling, train/test split, and Random Forest regression.'),
    ]

    for x, y, w, h, text_value in cards:
        shape = slide.shapes.add_shape(1, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = white
        shape.line.color.rgb = green
        shape.line.width = 1
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text_value
        p.font.size = Pt(17)
        p.font.color.rgb = dark

    metric1 = slide.shapes.add_shape(1, Inches(0.8), Inches(5.1), Inches(3.7), Inches(1.1))
    metric1.fill.solid(); metric1.fill.fore_color.rgb = green; metric1.line.fill.background()
    p = metric1.text_frame.paragraphs[0]
    p.text = 'R² Score: 0.9336'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = white
    p.alignment = PP_ALIGN.CENTER

    metric2 = slide.shapes.add_shape(1, Inches(4.9), Inches(5.1), Inches(3.7), Inches(1.1))
    metric2.fill.solid(); metric2.fill.fore_color.rgb = white; metric2.line.color.rgb = green
    p = metric2.text_frame.paragraphs[0]
    p.text = 'Rows Used: 1625'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = dark
    p.alignment = PP_ALIGN.CENTER

    metric3 = slide.shapes.add_shape(1, Inches(9.0), Inches(5.1), Inches(3.5), Inches(1.1))
    metric3.fill.solid(); metric3.fill.fore_color.rgb = green; metric3.line.fill.background()
    p = metric3.text_frame.paragraphs[0]
    p.text = 'Model: Random Forest Regressor'
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = white
    p.alignment = PP_ALIGN.CENTER


def add_dataset_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = light

    banner = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.8))
    banner.fill.solid(); banner.fill.fore_color.rgb = green; banner.line.fill.background()
    header = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(8), Inches(0.4))
    p = header.text_frame.paragraphs[0]
    p.text = 'Dataset and Feature Engineering'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = white

    left = slide.shapes.add_shape(1, Inches(0.6), Inches(1.2), Inches(5.8), Inches(5.0))
    left.fill.solid(); left.fill.fore_color.rgb = white; left.line.color.rgb = green; left.line.width = 1
    p = left.text_frame.paragraphs[0]
    p.text = 'Key Features\n\n• Latitude and longitude\n• NDVI, GNDVI, NDWI, SAVI\n• Soil moisture\n• Temperature\n• Rainfall\n• Crop type\n• Yield value'
    p.font.size = Pt(20)
    p.font.color.rgb = dark

    right = slide.shapes.add_shape(1, Inches(6.8), Inches(1.2), Inches(5.9), Inches(5.0))
    right.fill.solid(); right.fill.fore_color.rgb = white; right.line.color.rgb = green; right.line.width = 1
    p = right.text_frame.paragraphs[0]
    p.text = 'Engineering Added\n\n• NDVI × Temperature\n• NDVI × Rainfall\n• SAVI × Soil Moisture\n\nThese interactions capture crop stress, plant vigor, and moisture conditions more realistically.'
    p.font.size = Pt(19)
    p.font.color.rgb = dark


def add_workflow_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = light

    banner = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.8))
    banner.fill.solid(); banner.fill.fore_color.rgb = green; banner.line.fill.background()
    header = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(8.5), Inches(0.4))
    p = header.text_frame.paragraphs[0]
    p.text = 'Model Workflow and Deployment'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = white

    steps = slide.shapes.add_shape(1, Inches(0.7), Inches(1.5), Inches(12.0), Inches(4.8))
    steps.fill.solid(); steps.fill.fore_color.rgb = white; steps.line.color.rgb = green; steps.line.width = 1
    p = steps.text_frame.paragraphs[0]
    p.text = '1. Load and clean agricultural dataset\n2. Engineer interaction features\n3. Encode crop type\n4. Scale numerical inputs\n5. Split into training and testing sets\n6. Train Random Forest Regressor\n7. Evaluate using R² score\n8. Deploy web application for real-time prediction'
    p.font.size = Pt(22)
    p.font.color.rgb = dark

    lower = slide.shapes.add_shape(1, Inches(0.7), Inches(6.3), Inches(12.0), Inches(0.9))
    lower.fill.solid(); lower.fill.fore_color.rgb = green; lower.line.fill.background()
    p = lower.text_frame.paragraphs[0]
    p.text = 'Web App: user inputs field conditions and receives a predicted yield instantly'
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = white
    p.alignment = PP_ALIGN.CENTER


def add_results_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = light

    banner = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.8))
    banner.fill.solid(); banner.fill.fore_color.rgb = green; banner.line.fill.background()
    header = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(8), Inches(0.4))
    p = header.text_frame.paragraphs[0]
    p.text = 'Results and Impact'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = white

    left = slide.shapes.add_shape(1, Inches(0.8), Inches(1.5), Inches(4.5), Inches(4.8))
    left.fill.solid(); left.fill.fore_color.rgb = white; left.line.color.rgb = green; left.line.width = 1
    p = left.text_frame.paragraphs[0]
    p.text = 'Model Performance\n\n• R² Score: 0.9336\n• Strong explanation of yield variability\n• Suitable for planning and decision support'
    p.font.size = Pt(22)
    p.font.color.rgb = dark

    right = slide.shapes.add_shape(1, Inches(5.9), Inches(1.5), Inches(6.5), Inches(4.8))
    right.fill.solid(); right.fill.fore_color.rgb = white; right.line.color.rgb = green; right.line.width = 1
    p = right.text_frame.paragraphs[0]
    p.text = 'Business Value\n\n• Helps farmers plan cultivation\n• Supports policy and food security strategies\n• Enables agribusiness forecasting\n• Improves resource allocation and risk management'
    p.font.size = Pt(22)
    p.font.color.rgb = dark


def add_closing_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = green

    box = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.0), Inches(1.2))
    p = box.text_frame.paragraphs[0]
    p.text = 'Smart Agriculture Through Predictive Analytics'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = white
    p.alignment = PP_ALIGN.CENTER

    sub = slide.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(10.3), Inches(2.4))
    p = sub.text_frame.paragraphs[0]
    p.text = 'This project demonstrates how machine learning, field data, and remote sensing can be combined to deliver practical agricultural decision support and improved yield planning.'
    p.font.size = Pt(22)
    p.font.color.rgb = white
    p.alignment = PP_ALIGN.CENTER

    foot = slide.shapes.add_textbox(Inches(1.0), Inches(6.2), Inches(11.2), Inches(0.5))
    p = foot.text_frame.paragraphs[0]
    p.text = 'Thank You'
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = white
    p.alignment = PP_ALIGN.CENTER


add_title_slide('Crop Yield Prediction for Smart Agriculture', 'Inceptor Kenya | Data Science Capstone Presentation')
add_dataset_slide()
add_workflow_slide()
add_results_slide()
add_closing_slide()

output = 'Crop_Yield_Prediction_Presentation.pptx'
prs.save(output)
print(f'Created: {output}')
